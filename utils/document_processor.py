import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from typing import List, Dict

class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
            return ""

    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        try:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except Exception as e:
            print(f"Error reading DOCX {file_path}: {e}")
            return ""

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
            return ""

    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading TXT {file_path}: {e}")
            return ""

    def process_document(self, file_path: str) -> Dict:
        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_name)[1].lower()
        if file_extension == '.pdf':
            text = self.extract_text_from_pdf(file_path)
        elif file_extension == '.docx':
            text = self.extract_text_from_docx(file_path)
        elif file_extension == '.pptx':
            text = self.extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            text = self.extract_text_from_txt(file_path)
        else:
            print(f"Unsupported file type: {file_extension}")
            return None
        if not text:
            print(f"No text could be extracted from {file_name}")
            return None
        return {
            "text": text,
            "metadata": {"source": file_name}
        }
